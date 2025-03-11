import logging
import logging.handlers
import math
import os
import re
import subprocess
import sys
from threading import Thread

import cv2
import pandas as pd
import pydub
import PyPDF2
import requests
from _uix import UIXBot
from cli_colors import (
    BLUE,
    BWHITE,
    CYAN,
    DBLUE,
    DCYAN,
    DGREEN,
    DMAGENTA,
    DRED,
    DYELLOW,
    FCYAN,
    FMAGENTA,
    GREEN,
    ICYAN,
    IGREEN,
    MAGENTA,
    RED,
    RESET,
    YELLOW,
)
from colors import (
    Black,
    Blue,
    Cyan,
    Gray,
    Green,
    Light_Gray,
    Magenta,
    Orange,
    Purple,
    Red,
    Semi_Transparent_White,
    Transparent,
    White,
    Yellow,
)
from docx import Document
from formats import (
    SUPPORTED_AUDIO_FORMATS_INPUT,
    SUPPORTED_AUDIO_FORMATS_OUTPUT,
    SUPPORTED_IMAGE_FORMATS,
    SUPPORTED_VIDEO_FORMATS,
    Video_codecs,
)
from gtts import gTTS
from kivy.clock import Clock
from kivy.graphics import Canvas, Color, Rectangle, RoundedRectangle
from kivy.uix.floatlayout import FloatLayout
from kivy.uix.gridlayout import GridLayout
from kivy.uix.label import Label
from m4a_converter import _m4a_main_
from moviepy import VideoFileClip
from openpyxl import load_workbook
from pdf2docx import parse
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation

# from pydub import AudioSegment
from reportlab.lib.pagesizes import letter
from reportlab.platypus import Paragraph, SimpleDocTemplate
from rich.progress import Progress
from tqdm import tqdm
from kivymd.uix.label import MDLabel

_ext_word = ["doc", "docx"]
_ext_ppt_ = ["ppt", "pptx"]
_ext_xls = ["xls", "xlsx"]

logging.basicConfig(level=logging.INFO, format="%(levelname)-8s %(message)s")
logger = logging.getLogger(__name__)


class FileConverter:
    def __init__(self, callback, callback2, **kwargs):
        self.callback = callback
        self.callback2 = callback2

    # .....................Conversions........................
    def convert_xls_to_csv(self, obj):
        print(f"Called with file {obj}")

        def worker(file):
            print(file)
            try:
                read_engine_store = {
                    "xls": "xlrd",
                    "xlsx": "openpyxl",
                    "xlsb": "pyxlsb",
                    "odf": "ods",
                }
                """Load the Excel file"""
                df = pd.read_excel(
                    file, engine=read_engine_store.get(file.split(".")[-1].lower())
                )

                Clock.schedule_once(
                    lambda dt: self.callback(
                        item=Label(
                            text=f"Converting {file} to csv",
                            size_hint_y=None,
                            color=Black,
                            height=40,
                        ),
                        log=True,
                    )
                )

                total_rows = df.shape[0]

                Clock.schedule_once(
                    lambda dt: self.callback(
                        item=Label(
                            text=f"Length {total_rows} rows",
                            size_hint_y=None,
                            color=Blue,
                            height=40,
                        ),
                        log=True,
                    )
                )

                df.to_csv(file.split(".")[0] + "csv", index=False)

                Clock.schedule_once(
                    lambda dt: self.callback(
                        item=Label(
                            text="Completed", size_hint_y=None, color=Green, height=40
                        ),
                        log=True,
                    )
                )

                line = "." * 100

                Clock.schedule_once(
                    lambda dt: self.callback(
                        item=Label(text=line, color=Green, size_hint_y=None, height=40),
                        log=True,
                    )
                )

            except KeyboardInterrupt:
                print("\nQuit❕")
                exit(1)
            except Exception as e:
                Clock.schedule_once(
                    lambda dt: self.callback(
                        item=Label(text=str(e), padding=5, color=Red), error=True
                    )
                )
                # Clock.schedule_once(lambda dt: self._force_ui_update(), 0.5)
                # print(e)

        if os.path.isdir(obj):
            for _, file in enumerate(os.listdir(obj)):
                if os.path.isfile(file) and file.endswith(
                    ext for ext in ("xls", "xlsx")
                ):
                    Thread(name=_, target=worker(file))
        else:
            worker(obj)

    def convert_xls_to_word(self, excel):
        word_file = excel.split(".")[0] + ".docx"
        try:
            """Read the XLS file using pandas"""

            df = pd.read_excel(excel)

            """Create a new Word document"""
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"Converting {excel}..",
                        size_hint_y=None,
                        color=Cyan,
                        height=40,
                    ),
                    log=True,
                )
            )
            doc = Document()

            """Iterate over the rows of the dataframe and add them to the
            Word document"""
            logger.info(f"{ICYAN}Converting {excel}..{RESET}")

            # time.sleep(2)
            Clock.schedule_once(lambda dt: self.callback2())
            total_rows = df.shape[0]
            for _, row in df.iterrows():
                current_row = _ + 1
                percentage = (current_row / total_rows) * 100
                for value in row:
                    doc.add_paragraph(str(value))
                Clock.schedule_once(
                    lambda _: self.callback(
                        item=None,
                        current=current_row,
                        _max=total_rows,
                        update=True,
                        progress=True,
                    )
                )
                print(
                    f"Row {DYELLOW}{current_row}/{total_rows} {DBLUE}{percentage:.1f}%{RESET}",
                    end="\r",
                )
                # print(f"\033[1;36m{row}{RESET}")

            # Save the Word document
            doc.save(word_file)
            print(f"{DGREEN}Conversion successful!{RESET}", end="\n")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="Conversion successful!",
                        size_hint_y=None,
                        color=Green,
                        height=40,
                    ),
                    log=True,
                )
            )
        except KeyboardInterrupt:
            print("\nQuit⌨️")
            sys.exit(1)
        except Exception as e:
            err = str(e)
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(text=err, size_hint_y=None, color=Red, height=40),
                    error=True,
                )
            )
            print(f"{RED}Oops Conversion failed:❕{RESET}", str(e))

    def convert_csv_to_xls(self, file):
        csv_list = [file]
        csv_list = [item for item in csv_list if item.split(".")[-1].lower() in ("csv")]
        try:
            with Progress() as progress:
                task = progress.add_task("[cyan]Coverting", total=len(csv_list))
                for file in csv_list:
                    file_name = file[:-3] + "xlsx"
                    df = pd.read_csv(file)
                    # excel engines ('openpyxl' or 'xlsxwriter')
                    df.to_excel(file_name, engine="openpyxl", index=False)

                    # Load the workbook and the sheet
                    workbook = load_workbook(file_name)
                    sheet = workbook.active

                    Clock.schedule_once(lambda dt: self.callback2())

                    # Convert columns generator to a list and get the length
                    total = len(list(sheet.columns))
                    for current, column in enumerate(sheet.columns):
                        Clock.schedule_once(
                            lambda dt: self.callback(
                                item=None,
                                current=current + 1,
                                _max=total,
                                update=True,
                                progress=True,
                            )
                        )
                        max_length = 0
                        column_letter = column[0].column_letter

                        max_length = max(len(str(cell.value)) for cell in column)
                        """for cell in column:
                            try:
                                if len(str(cell.value)) > max_length:
                                    max_length = len(cell.value)

                            except Exception:
                                pass
                        """
                        adjusted_width = max_length + 2
                        sheet.column_dimensions[column_letter].width = adjusted_width

                    # Save the workbook
                    workbook.save(file_name)
                    progress.update(task, advance=1)
        except Exception as e:
            err = str(e)
            print(f"{RED}{e}{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(text=err, size_hint_y=None, color=Red, height=40),
                    error=True,
                )
            )

    def convert_word_to_text(self, file):
        try:
            logger.info(f"{DYELLOW}Create Doument Tablet{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="Create Doument Tablet",
                        size_hint_y=None,
                        color=Yellow,
                        height=40,
                    ),
                    log=True,
                )
            )
            doc = Document(file)

            Clock.schedule_once(lambda dt: self.callback2())
            total = len(doc.paragraphs)
            with open(file, "w", encoding="utf-8") as f:
                for Par, paragraph in enumerate(doc.paragraphs):
                    Clock.schedule_once(
                        lambda dt: self.callback(
                            item=None,
                            current=Par + 1,
                            _max=total,
                            update=True,
                            progress=True,
                        )
                    )
                    f.write(paragraph.text + "\n")

                    print(f"Par:{BLUE}{Par+1}/{total}{RESET}", end="\r")
                logger.info(f"{DMAGENTA}Conversion of file to txt success{RESET}")
                Clock.schedule_once(
                    lambda dt: self.callback(
                        item=Label(
                            text="[*] Conversion of file to txt success",
                            size_hint_y=None,
                            color=Green,
                            height=40,
                        ),
                        log=True,
                    )
                )

            return file
        except KeyboardInterrupt:
            print("\nQuit❕⌨️")
            sys.exit()
        except Exception as e:
            err = str(e)
            logger.error(e)
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=err, size_hint_y=None, bold=True, color=Red, height=40
                    ),
                    error=True,
                )
            )

    def convert_pdf_to_text(self, file):
        txt_file = file[:-3] + "txt"
        try:
            print(f"{DYELLOW}Open and read the pdf document..{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="Open and read the pdf document..",
                        size_hint_y=None,
                        color=Cyan,
                        height=40,
                    ),
                    log=True,
                )
            )
            with open(file, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""

                print(f"{YELLOW}Convert pages..{RESET}")
                Clock.schedule_once(
                    lambda dt: self.callback(
                        item=Label(
                            text="Convert pages",
                            size_hint_y=None,
                            color=Magenta,
                            height=40,
                        ),
                        log=True,
                    )
                )

                Clock.schedule_once(lambda dt: self.callback2())
                _max = len(pdf_reader.pages)
                for page_num in range(_max):
                    Clock.schedule_once(
                        lambda dt: self.callback(
                            item=None,
                            current=page_num + 1,
                            _max=_max,
                            update=True,
                            progress=True,
                        )
                    )
                    logger.info(
                        f"Page {DBLUE}{page_num+1}{RESET}/{len(pdf_reader.pages)}"
                    )
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text()

            with open(txt_file, "w", encoding="utf-8") as f:
                f.write(text)

            logger.info(f"{MAGENTA}New file is {CYAN}{txt_file}{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"[+] New file is {txt_file}",
                        size_hint_y=None,
                        color=Purple,
                        bold=True,
                        height=40,
                    ),
                    log=True,
                )
            )
            logger.info(f"{DGREEN}Success👨‍💻✅{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"[✅] Success👨‍💻✅",
                        size_hint_y=None,
                        color=Yellow,
                        height=40,
                    ),
                    log=True,
                )
            )
        except Exception as e:
            err = str(e)
            logger.error(f"{RED}{e}{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(text=err, size_hint_y=None, color=Red, height=40),
                    error=True,
                )
            )

    def convert_pdf_to_word(self, file):
        word_file = file.split(".")[0] + "docx"

        try:
            print(f"{DYELLOW}Parse the pdf document..{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="Parse the pdf document..",
                        size_hint_y=None,
                        color=Yellow,
                        height=40,
                    ),
                    log=True,
                )
            )

            parse(file, word_file, start=0, end=None)

            logger.info(f"{MAGENTA}New file is {CYAN}{word_file}{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"[+] New file is {word_file}",
                        size_hint_y=None,
                        color=Purple,
                        bold=True,
                        height=40,
                    ),
                    log=True,
                )
            )
            logger.info(f"{DGREEN}Success👨‍💻✅{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"[✅] Success👨‍💻✅",
                        size_hint_y=None,
                        color=Yellow,
                        height=40,
                    ),
                    log=True,
                )
            )
        except KeyboardInterrupt:
            print("\nQuit❕")
            sys.exit(1)
        except Exception as e:
            err = str(e)
            logger.info(f"{DRED}All conversion attempts have failed: {e}{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(text=err, size_hint_y=None, color=Red, height=40),
                    error=True,
                )
            )

    def txt_to_pdf(self, file):
        _pdf_ = file.split(".")[0] + ".pdf"
        # Read the contents of the input .txt file
        with open(file, "r", encoding="utf-8") as file:
            text_contents = file.readlines()

        # Initialize the PDF document
        logger.info(f"{DYELLOW}Initialize the PDF document{RESET}")
        Clock.schedule_once(
            lambda dt: self.callback(
                item=Label(
                    text="Initialize the PDF document",
                    size_hint_y=None,
                    color=Yellow,
                    height=40,
                ),
                log=True,
            )
        )
        doc = SimpleDocTemplate(_pdf_, pagesize=letter)

        # Create a story to hold the elements of the PDF
        logger.info(f"{DYELLOW}Create a story to hold the c of the PDF{RESET}")
        Clock.schedule_once(
            lambda dt: self.callback(
                item=Label(
                    text="Create a story to hold the content of the PDF",
                    size_hint_y=None,
                    color=Yellow,
                    height=40,
                ),
                log=True,
            )
        )
        story = []

        # Iterate through each line in the input .txt file and add it to the PDF
        logger.info(f"{DYELLOW}Add line to PDF file{RESET}")
        Clock.schedule_once(
            lambda dt: self.callback(
                item=Label(
                    text="Add lines to PDF file",
                    size_hint_y=None,
                    color=Yellow,
                    height=40,
                ),
                log=True,
            )
        )
        try:
            Clock.schedule_once(lambda dt: self.callback2())
            length = len(text_contents)
            for _line_count_, line in enumerate(text_contents):
                logger.info(f"Lines {DBLUE}{_line_count_+1}{RESET}/{length}")
                Clock.schedule_once(
                    lambda dt: self.callback(
                        item=None,
                        current=_line_count_ + 1,
                        _max=length,
                        update=True,
                        progress=True,
                    )
                )
                story.append(Paragraph(line.strip(), style="normalText"))

        except KeyboardInterrupt:
            print("\nQuit❕⌨️")
            sys.exit(1)
        except Exception as e:
            err = str(e)
            logger.error(e)
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(text=err, size_hint_y=None, color=Red, height=40),
                    error=True,
                )
            )
        finally:
            # Build and write the PDF document
            logger.info(f"{DYELLOW}Build and write the PDF document{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="Build and write the PDF document",
                        size_hint_y=None,
                        color=Yellow,
                        height=40,
                    ),
                    log=True,
                )
            )
            doc.build(story)
            logger.info(f"{MAGENTA}New file is {CYAN}{_pdf_}{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"New file is:> {_pdf_}",
                        size_hint_y=None,
                        color=Cyan,
                        height=40,
                    ),
                    log=True,
                )
            )
            print(f"\n{DGREEN}Success👨‍💻✅{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="Success 👨‍💻✅",
                        size_hint_y=None,
                        bold=True,
                        color=Green,
                        height=40,
                    ),
                    log=True,
                )
            )

    def word_to_pptx(self, file):
        try:
            # Load the Word document
            print(f"{DYELLOW}Load the Word document..{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="Load the Word document..",
                        size_hint_y=None,
                        color=Yellow,
                        height=40,
                    ),
                    log=True,
                )
            )
            doc = Document(file)
            pptx = file.split(".")[0] + ".pptx"

            # Create a new PowerPoint presentation
            print(f"{DYELLOW}Create a new PowerPoint presentation..{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=MDLabel(
                        text="Create a new [b]PowerPoint[/b] presentation..",
                        size_hint_y=None,
                        text_color=Yellow,
                        halign="center",
                        markup=True,
                        height=40,
                    ),
                    log=True,
                )
            )
            prs = Presentation()

            length = len(doc.paragraphs)
            # Iterate through each paragraph in the Word document
            print(
                f"{DGREEN}Populating pptx slides with {DYELLOW}{length}{DGREEN} entries..{RESET}"
            )
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"[+] Populating pptx slides with {length}",
                        size_hint_y=None,
                        color=Magenta,
                        height=40,
                    ),
                    log=True,
                )
            )

            Clock.schedule_once(lambda dt: self.callback2())
            for count, paragraph in enumerate(doc.paragraphs):
                perc = count + 1 / length
                perc *= 100
                print(f"{DMAGENTA}Progress:: {DCYAN}{perc:.2f}%{RESET}", end="\r")
                Clock.schedule_once(
                    lambda dt: self.callback(
                        item=None,
                        current=count + 1,
                        _max=length,
                        update=True,
                        progress=True,
                    )
                )
                # Create a new slide in the PowerPoint presentation
                slide = prs.slides.add_slide(prs.slide_layouts[1])

                # Add the paragraph text to the slide
                slide.shapes.title.text = paragraph.text

            # Save the PowerPoint presentation
            prs.save(pptx)
            logger.info(f"{MAGENTA}New file is {CYAN}{pptx}{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"[+] New file is {pptx}",
                        size_hint_y=None,
                        color=Purple,
                        bold=True,
                        height=40,
                    ),
                    log=True,
                )
            )
            logger.info(f"{DGREEN}Success👨‍💻✅{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"[✅] Success👨‍💻✅",
                        size_hint_y=None,
                        color=Yellow,
                        height=40,
                    ),
                    log=True,
                )
            )
        except KeyboardInterrupt:
            print("\nQuit❕⌨️")
            sys.exit(1)
        except Exception as e:
            err = str(e)
            logger.error(e)
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(text=err, size_hint_y=None, color=Red, height=40),
                    error=True,
                )
            )

    def text_to_word(self, file):
        word_file = file.split(".")[0] + "docx"
        try:
            # Read the text file
            logger.info(f"{DCYAN}Open and read the text file{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="Open and read the text file",
                        size_hint_y=None,
                        color=Cyan,
                        height=40,
                    ),
                    log=True,
                )
            )
            with open(file, "r", encoding="utf-8", errors="ignore") as file:
                text_content = file.read()

            # Filter out non-XML characters
            filtered_content = re.sub(
                r"[^\x09\x0A\x0D\x20-\uD7FF\uE000-\uFFFD]+", "", text_content
            )

            # Create a new Word document
            logger.info(f"{DYELLOW}Create Doument Tablet{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="Create Doument Tablet",
                        size_hint_y=None,
                        color=Yellow,
                        height=40,
                    ),
                    log=True,
                )
            )
            doc = Document()
            # Add the filtered text content to the document
            doc.add_paragraph(filtered_content)

            # Save the document as a Word file
            doc.save(word_file)
            logger.info(f"{MAGENTA}New file is {DCYAN}{word_file}{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"[+] New file is {word_file}",
                        size_hint_y=None,
                        color=Purple,
                        bold=True,
                        height=40,
                    ),
                    log=True,
                )
            )
            logger.info(f"{DGREEN}Success👨‍💻✅{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"[✅] Success👨‍💻✅",
                        size_hint_y=None,
                        color=Yellow,
                        height=40,
                    ),
                    log=True,
                )
            )
        except FileExistsError as e:
            err = str(e)
            logger.error(f"{str(e)}📁")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"[📁] {err}", size_hint_y=None, color=Yellow, height=40
                    ),
                    log=True,
                )
            )
        except Exception as e:
            err = str(e)
            logger.error(f"{RED}{e}{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(text=err, size_hint_y=None, color=Red, height=40),
                    error=True,
                )
            )

    def convert_word_to_pdf(self, word_file):
        pdf_file_dir = os.path.dirname(word_file)
        pdf_file = os.path.splitext(word_file)[0] + ".pdf"

        try:
            print(f"{BLUE}Converting: {RESET}{word_file} {BLUE}to {RESET}{pdf_file}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text=f"[+] Converting: {word_file}",
                        size_hint_y=None,
                        bold=True,
                        underline=True,
                        color=Yellow,
                        height=40,
                    ),
                    log=True,
                )
            )
            if os.name == "posix":  # Check if running on Linux
                # Use subprocess to run the dpkg and grep commands
                result = subprocess.run(
                    ["dpkg", "-l", "libreoffice"], stdout=subprocess.PIPE, text=True
                )
                if result.returncode != 0:
                    logger.exception(f"{RED}Libreoffice not found !{RESET}")
                    Clock.schedule_once(
                        lambda dt: self.callback(
                            item=Label(
                                text=f"Libreoffice not found !",
                                size_hint_y=None,
                                color=Red,
                                height=40,
                            ),
                            error=True,
                        )
                    )
                    print(f"{CYAN}Initiating critical redundacy measure !{RESET}")
                    Clock.schedule_once(
                        lambda dt: self.callback(
                            item=Label(
                                text="Initiating critical redundacy measure !",
                                size_hint_y=None,
                                color=Yellow,
                                height=40,
                            ),
                            log=True,
                        )
                    )
                    self.word2pdf_extra(word_file)
                subprocess.run(
                    [
                        "soffice",
                        "--convert-to",
                        "pdf",
                        word_file,
                        "--outdir",
                        pdf_file_dir,
                    ]
                )

                print(
                    f"{DMAGENTA} Successfully converted {word_file} to {pdf_file}{RESET}"
                )
                Clock.schedule_once(
                    lambda dt: self.callback(
                        item=Label(
                            text="Done ✅",
                            size_hint_y=None,
                            color=Green,
                            bold=True,
                            height=40,
                        ),
                        log=True,
                    )
                )
                return pdf_file

            elif os.name == "nt":
                self.word2pdf_extra(word_file)
                return pdf_file

        except Exception as e:
            err = str(e)
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(text=err, size_hint_y=None, color=Red, height=40),
                    error=True,
                )
            )
            print(f"Error converting {word_file} to {pdf_file}: {e}")

    def word2pdf_extra(self, obj, outf=None):
        """For window users since it requires Microsoft word to be installed"""
        if obj.split(".")[-1] not in {"doc", "docx"}:
            print(f"{RED}File is not a word file{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="File is not a word file",
                        size_hint_y=None,
                        color=Purple,
                        height=40,
                    ),
                    error=True,
                )
            )
            sys.exit(1)
        pdf_file = os.path.splitext(obj)[0] + ".pdf" if outf is None else outf
        try:
            from docx2pdf import convert

            convert(obj, pdf_file)
            print(f"{GREEN}Conversion ✅{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="Conversion ✅", size_hint_y=None, color=Green, height=40
                    ),
                    log=True,
                )
            )
            sys.exit(0)
        except ImportError:
            print(f"{RED}docx2pdf Not found. {CYAN}Run pip install docx2pdf{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(
                        text="[-] Import error :- docx2pdf.",
                        size_hint_y=None,
                        color=Red,
                        height=40,
                    ),
                    error=True,
                )
            )
        except Exception as e:
            err = str(e)
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(text=err, size_hint_y=None, color=Red, height=40),
                    error=True,
                )
            )
            logger.error(e)

    def pptx_to_txt(self, file_path, dest=None):
        """
        Converts a PowerPoint presentation to a text file, extracting text from slides.

        Args:
            file_path (str): The path to the PowerPoint presentation file (.ppt or .pptx).
            dest (str, optional): Destination type, if "text", returns text buffer instead of file path. Defaults to None.

        Returns:
            str: The path to the generated text file (.txt), or text buffer if dest="text", or None in case of error.
        """
        txt_file = self._initialize_txt_file_path(file_path)

        try:
            file_path = os.path.abspath(file_path)
            file_path = self._convert_to_pptx_if_ppt(file_path)

            presentation = Presentation(file_path)
            total_slides = len(presentation.slides)

            self._log_and_display_slide_count(total_slides)
            Clock.schedule_once(
                lambda dt: self.callback2()
            )  # What is callback2 for? Clarify or rename.

            with Progress() as progress:
                task_id = progress.add_task("[magenta]Preparing..", total=total_slides)
                for count, slide in enumerate(presentation.slides, 1):
                    self._update_progress(count, total_slides, progress, task_id)
                    self._extract_slide_text_to_file(slide, txt_file)
                    progress.update(task_id, advance=1)

            return self._handle_destination_output(dest, txt_file)

        except Exception as e:
            return self._handle_exception(e, detailed_error_prefix="\n❌Oops! ")

    def _initialize_txt_file_path(self, file_path):
        """
        Initializes the text output file path based on the PowerPoint file path.
        """
        return (os.path.splitext(file_path)[0]) + ".txt"

    def _convert_to_pptx_if_ppt(self, file_path):
        """
        Converts a PPT file to PPTX if the input file is in PPT format.
        """
        ext = os.path.splitext(file_path)[-1][1:]
        if ext == "ppt":
            return self.convert_ppt_to_pptx(file_path)
        return file_path

    def _log_and_display_slide_count(self, slide_count):
        """
        Logs the slide count and displays it using a callback.
        """
        message = f"Target images := {slide_count}"
        logger.info(f"Slide count ={DMAGENTA} {slide_count}{RESET}")
        self._log_and_display_message(message, Yellow)

    def _update_progress(self, current, maximum, progress_obj, task_id):
        """
        Updates the progress bar and display.
        """
        Clock.schedule_once(
            lambda dt: self.callback(
                item=None,
                current=current,
                _max=maximum,
                update=True,
                progress=True,
            )
        )

    def _extract_slide_text_to_file(self, slide, txt_file):
        """
        Extracts text from each shape in a PowerPoint slide and writes it to the text file.
        """
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    self._process_paragraph_text_to_file(paragraph, txt_file)

    def _process_paragraph_text_to_file(self, paragraph, txt_file):
        """
        Processes a paragraph from PowerPoint, extracting and writing text runs to the text file.
        """
        if any(run.text.strip() for run in paragraph.runs):
            for run in paragraph.runs:
                text = run.text.strip()
                if text and text != " ":
                    self._write_text_to_file(text, txt_file)

    def _write_text_to_file(self, text, txt_file):
        """
        Writes text to the specified text file.
        """
        with open(txt_file, "a") as fl:
            fl.write(text)

    def _handle_destination_output(self, dest, txt_file):
        """
        Handles the final output based on the destination type.
        """
        if dest == "text":
            return self._read_text_buffer(txt_file)
        else:
            self._log_and_display_new_file_created(txt_file)
            self._log_and_display_success()
            return txt_file

    def _read_text_buffer(self, txt_file):
        """
        Reads the content of the text file into a text buffer and returns it.
        """
        with open(txt_file, "r") as fl:
            text_buffer = fl.read()
        return text_buffer

    def _log_and_display_new_file_created(self, txt_file):
        """
        Logs and displays a message indicating a new file has been created.
        """
        message = f"New file := {txt_file}"
        logger.info(f"{MAGENTA}New file is {CYAN}{txt_file}{RESET}")
        self._log_and_display_message(message, Magenta)

    def _log_and_display_success(self):
        """
        Logs and displays a success message.
        """
        logger.info(f"{DGREEN}Success👨‍💻✅{RESET}")
        self._log_and_display_message("Success", Green)

    def _handle_exception(self, e, detailed_error_prefix=""):
        """
        Handles exceptions during the PowerPoint to text conversion process.
        """
        err_msg = str(e)
        logger.error(f"{detailed_error_prefix}{RED}{e}{RESET}")
        self._log_and_display_error_message(err_msg)
        return None

    def _log_and_display_error_message(self, error_message):
        """
        Logs and displays an error message using a callback.
        """
        Clock.schedule_once(
            lambda dt: self.callback(
                item=Label(text=error_message, size_hint_y=None, color=Red, height=40),
                log=True,  # Keep log as True for error logging
            )
        )

    def convert_ppt_to_pptx(self, obj: os.PathLike):
        import platform

        try:
            if obj.endswith("ppt"):
                if platform.system() in ("Linux", "MacOS") or os.name == "posix":
                    subprocess.run(
                        ["soffice", "--headless", "--convert-to", "pptx", obj]
                    )
                    return os.path.splitext(obj)[0] + ".pptx"
                elif platform.system() in ("Windows") or os.name == "nt":
                    import win32com.client

                    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                    powerpoint.Visible = 1
                    ppt = powerpoint.Presentations.Open(obj)
                    pptx_file = os.path.splitext(obj)[0] + ".pptx"
                    ppt.SaveAs(pptx_file, 24)  # 24 is the format for pptx
                    ppt.Close()
                    powerpoint.Quit()
                    return pptx_file
            else:
                print(f"{RED}Unable to identify the system{RESET}")
                Clock.schedule_once(
                    lambda dt: self.callback(
                        item=Label(
                            text="Unable to identify the system",
                            size_hint_y=None,
                            color=Red,
                            height=40,
                        ),
                        err=True,
                    )
                )
        except KeyboardInterrupt:
            print("\nQuit!")
            exit(1)
        except Exception as e:
            err = str(e)
            logger.error(f"{RED}{e}{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(text=err, size_hint_y=None, color=Red, height=40),
                    error=True,
                )
            )

    def ppt_to_word(self, file_path):
        """
        Converts a PowerPoint presentation to a Word document, preserving text formatting.

        Args:
            file_path (str): The path to the PowerPoint presentation file (.ppt or .pptx).

        Returns:
            str: The path to the generated Word document file (.docx), or None in case of error.
        """
        word_file = self._initialize_word_file_path(file_path)

        try:
            self._log_and_display_message("Create Doument Tablet", Yellow)
            file_path = os.path.abspath(file_path)
            file_path = self._convert_to_pptx_if_ppt(file_path)

            presentation = Presentation(file_path)
            document = Document()
            total_slides = len(presentation.slides)

            self._log_and_display_message(f"Slide count ={total_slides}", Magenta)

            Clock.schedule_once(
                lambda dt: self.callback2()
            )  # What is callback2 for? Consider renaming or clarifying.

            with Progress() as progress:
                task_id = progress.add_task("[magenta]Preparing..", total=total_slides)
                for count, slide in enumerate(
                    presentation.slides, 1
                ):  # start enumerate from 1 for count clarity
                    self._update_progress(count, total_slides, progress, task_id)
                    self._process_slide_text(slide, document)
                    progress.update(task_id, advance=1)

            document.save(word_file)
            self._log_and_display_message(f"New file := {word_file}", Magenta)
            self._log_and_display_message("Success👨‍💻✅", Green)
            return word_file

        except Exception as e:
            return self._handle_exception(e)

    def _initialize_word_file_path(self, file_path):
        """
        Initializes the Word output file path based on the PowerPoint file path.
        """
        return file_path.split(".")[0] + ".docx"

    def _log_and_display_message(self, message, color, log=True, error=False):
        """
        Logs a message and displays it using a callback mechanism.
        """
        logger.info(
            f"{self._get_color_code(color)}{message}{RESET}"
        )  # use a helper to get color code
        Clock.schedule_once(
            lambda dt: self.callback(
                item=Label(text=message, size_hint_y=None, color=color, height=40),
                log=log,
                error=error,
            )
        )

    def _get_color_code(self, color):
        """
        Helper function to get color escape codes based on color name.
        """
        color_map = {
            Yellow: DYELLOW,
            Magenta: DMAGENTA,
            Green: DGREEN,
            Red: RED,
            Cyan: CYAN,
        }
        return color_map.get(color, "")  # default to empty string if color not found

    def _process_slide_text(self, slide, document):
        """
        Extracts and processes text from each shape in a PowerPoint slide,
        adding formatted paragraphs to the Word document.
        """
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    self._process_paragraph_text(paragraph, document)

    def _process_paragraph_text(self, paragraph, document):
        """
        Processes a paragraph from PowerPoint, creating a corresponding paragraph in Word,
        preserving text formatting (bold, italic, underline, font, color).
        """
        if not any(run.text.strip() for run in paragraph.runs):
            return  # Skip empty paragraphs

        new_paragraph = document.add_paragraph()
        self._set_paragraph_formatting(new_paragraph)

        for run in paragraph.runs:
            if run.text.strip():
                self._process_run_text(run, new_paragraph)

    def _set_paragraph_formatting(self, paragraph):
        """
        Sets standard formatting for a Word paragraph.
        """
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
        from docx.shared import Pt

        paragraph.alignment = WD_PARAGRAPH_ALIGNMENT.JUSTIFY
        paragraph.space_after = Pt(6)
        paragraph.space_before = Pt(6)
        paragraph.line_spacing = 1.15

    def _process_run_text(self, run, paragraph):
        """
        Processes a text run from PowerPoint, adding a corresponding run in Word,
        preserving text formatting attributes.
        """
        from docx.shared import RGBColor as docxRGBColor
        from pptx.dml.color import RGBColor as pptxRGBColor

        new_run = paragraph.add_run(run.text)
        new_run.bold = run.font.bold
        new_run.italic = run.font.italic
        new_run.underline = run.font.underline
        new_run.font.name = run.font.name
        new_run.font.size = run.font.size

        try:
            if run.font.color and run.font.color.rgb:
                pptx_color = run.font.color.rgb
                docx_color = (
                    docxRGBColor(0, 0, 0)
                    if pptx_color == pptxRGBColor(255, 255, 255)
                    else docxRGBColor(*pptx_color)
                )
                new_run.font.color.rgb = docx_color
        except AttributeError:
            pass  # Handle cases where color attributes are not available

    def _handle_exception_simple(self, e):
        """
        Handles exceptions during the PowerPoint to Word conversion process.
        """
        err_msg = str(e)
        logger.error(f"{RED}{e}{RESET}")
        self._log_and_display_message(
            err_msg, Red, log=False, error=True
        )  # log=False to prevent double logging?
        return None  # or raise exception again if you want to propagate it

    def pdf2image(self, file, outf="png"):
        ###############################################################################
        """Create image objects from given files"""
        ###############################################################################
        outf = "png" if outf not in ("png", "jpg") else outf
        file_list = [file]
        imgs = []
        try:
            for file in tqdm(file_list):
                if file.lower().endswith("pdf"):
                    # Convert the PDF to a list of PIL image objects
                    print(f"{DBLUE}Generate image objects ..{RESET}")
                    Clock.schedule_once(
                        lambda dt: self.callback(
                            item=Label(
                                text="Generate image objects",
                                size_hint_y=None,
                                color=Blue,
                                height=40,
                            ),
                            log=True,
                        )
                    )
                    images = convert_from_path(file)

                    # Save each image to a file
                    length = len(images)
                    print(f"{YELLOW}Target images{BLUE} {len(images)}{RESET}")
                    Clock.schedule_once(
                        lambda dt: self.callback(
                            item=Label(
                                text=f"Target images := {length}",
                                size_hint_y=None,
                                color=Yellow,
                                height=40,
                            ),
                            log=True,
                        )
                    )

                    Clock.schedule_once(lambda dt: self.callback2())
                    for i, image in enumerate(images):
                        i += 1
                        Clock.schedule_once(
                            lambda _: self.callback(
                                item=None,
                                current=i,
                                _max=length,
                                update=True,
                                progress=True,
                            )
                        )
                        print(f"{DBLUE}{i}{RESET}", end="\r")
                        yd = f"{file.split('.')[0]}_{i+1}.{outf}"
                        image.save(yd)
                        imgs.append(yd)
                print(f"\n{GREEN}Ok{RESET}")

                return imgs
        except Exception as e:
            err = str(e)
            print(f"{RED}{e}{RESET}")
            Clock.schedule_once(
                lambda dt: self.callback(
                    item=Label(text=err, size_hint_y=None, color=Red, height=40),
                    error=True,
                )
            )


class ScannerBot:
    """
    Implementation of scanning to extract data from pdf files and images
    input_file -> file to be scanned pdf,image
    """

    def __init__(self, callback, **kwargs):
        self.callback = callback

    def preprocess(self):
        files_to_process = []

        if os.path.isfile(self.input_file):
            files_to_process.append(self.input_file)
        elif os.path.isdir(self.input_file):
            for file in os.listdir(self.input_file):
                file_path = os.path.join(self.input_file, file)
                if os.path.isfile(file_path):
                    files_to_process.append(file_path)

        return files_to_process

    def scanPDF(self, obj=None):
        """Obj- object for scanning where the object is not a list"""
        pdf_list = self.preprocess()
        pdf_list = [item for item in pdf_list if item.lower().endswith("pdf")]
        if obj:
            pdf_list = [obj]

        for pdf in tqdm(pdf_list):
            out_f = pdf[:-3] + "txt"
            print(f"{YELLOW}Read pdf ..{RESET}")

            with open(pdf, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                text = ""

                pg = 0
                for page_num in range(len(reader.pages)):
                    pg += 1

                    print(f"{DYELLOW}Progress:{RESET}", end="")
                    print(f"{CYAN}{pg}/{len(reader.pages)}{RESET}", end="\r")
                    page = reader.pages[page_num]
                    text += page.extract_text()

            print(f"\n{text}")
            print(f"\n{YELLOW}Write text to {GREEN}{out_f}{RESET}")
            with open(out_f, "w") as f:
                f.write(text)

            print(f"\n{DGREEN}Ok{RESET}")

    def scanAsImgs(self):
        file = self.input_file
        mc = FBot()
        img_objs = mc.pdf2image(file)
        # print(img_objs)
        from OCRTextExtractor import ExtractText

        text = ""
        for i in img_objs:
            extract = ExtractText(i)
            _text = extract.OCR()

            if _text is not None:
                text += _text
                with open(f"{self.input_file[:-4]}_filemac.txt", "a") as _writer:
                    _writer.write(text)

        def _cleaner_():
            print(f"{FMAGENTA}Clean")
            for obj in img_objs:
                if os.path.exists(obj):
                    print(obj, end="\r")
                    os.remove(obj)
                txt_file = f"{obj[:-4]}.txt"
                if os.path.exists(txt_file):
                    print(f"{FCYAN}{txt_file}{RESET}", end="\r")
                    os.remove(txt_file)

        # Do clean up
        _cleaner_()
        from overwrite import clear_screen

        clear_screen()
        print(f"{DGREEN}{IGREEN}Full Text{RESET}")
        print(text)
        print(f"{BWHITE}Text File ={IGREEN}{self.input_file[:-4]}_filemac.txt{RESET}")
        print(f"{GREEN}Ok✅{RESET}")
        return text

    def scanAsLongImg(self, file):
        """Convert the pdf to long image for scanning - text extraction"""
        file = self.input_file
        from longImg import LImage

        LI = LImage(file)
        fl = LI.preprocess()
        from OCRTextExtractor import ExtractText

        # fpath = file.split('.')[0] + '.png'
        tx = ExtractText(fl)
        text = tx.OCR()
        if text is not None:
            print(text)
            print(f"{GREEN}Ok{RESET}")
        return text


class TextToSpeechConverter:
    def __init__(self, callback, **kwargs):
        self.callback = callback

    from rich.progress import BarColumn, Progress, SpinnerColumn, TextColumn
    from kivy.logger import Logger as logger
    from kivy.utils import get_color_from_hex

    Yellow = get_color_from_hex("#F7E01C")

    def Synthesise(
        self,
        text: str,
        output_file: str,
        CHUNK_SIZE: int = 1_000,
        _tmp_folder_: str = "tmp_dir",
        max_retries: int = 30,
    ) -> None:
        """Converts given text to speech using Google Text-to-Speech API with improved structure."""

        (
            out_dir,
            _file_,
            _full_output_path_,
            checkpoint_file,
            start_chunk,
            resume_chunk_pos,
        ) = self._prepare_synthesis_environment(output_file, _tmp_folder_)
        total_chunks = math.ceil(len(text) / CHUNK_SIZE)
        counter = (
            math.ceil(resume_chunk_pos / CHUNK_SIZE) if resume_chunk_pos != 0 else 0
        )
        attempt = 0

        while attempt <= max_retries:
            try:
                with Progress(
                    SpinnerColumn(),
                    TextColumn("[green]Processing"),
                    BarColumn(),
                    TextColumn("[progress.percentage]{task.percentage:>3.1f}%"),
                ) as progress:
                    overall_task = progress.add_task(
                        "Overall Progress [magenta]Total", total=total_chunks
                    )
                    progress.update(overall_task, completed=counter)

                    for i in range(resume_chunk_pos, len(text), CHUNK_SIZE):
                        chunk = text[i : i + CHUNK_SIZE]
                        output_filename = self._get_output_filename_for_chunk(
                            _full_output_path_, counter, start_chunk
                        )

                        progress.console.print(
                            f"Processing chunk {counter + 1}/{total_chunks}\n",
                            end="\r",
                        )

                        self._synthesize_chunk_and_save(chunk, output_filename)
                        self._save_checkpoint(
                            checkpoint_file, counter + 1
                        )  # Increment counter before saving
                        counter += 1
                        progress.update(overall_task, advance=1)

                self._handle_synthesis_success(
                    _tmp_folder_, output_file, len(os.listdir(_tmp_folder_))
                )
                break  # Exit retry loop on success

            except requests.exceptions.ConnectionError:
                logger.error(f"{DRED}Connection Error{RESET}")
                resume_chunk_pos, attempt = self._handle_retry_attempt(
                    attempt, checkpoint_file
                )
            except requests.exceptions.HTTPError as e:
                logger.error(f"HTTP error: {e.status_code} - {e.reason}")
                resume_chunk_pos, attempt = self._handle_retry_attempt(
                    attempt, checkpoint_file
                )
            except requests.exceptions.RequestException as e:
                logger.error(f"{e}")
                resume_chunk_pos, attempt = self._handle_retry_attempt(
                    attempt, checkpoint_file
                )
            except (
                ConnectionError,
                ConnectionAbortedError,
                ConnectionRefusedError,
                ConnectionResetError,
            ):
                logger.error(
                    f"{RED}Connection at attempt {CYAN}{attempt+1}/{max_retries}{RESET}"
                )
                resume_chunk_pos, attempt = self._handle_retry_attempt(
                    attempt, checkpoint_file
                )
            except Exception as e:
                logger.error(f"{DMAGENTA}{attempt+1}/{max_retries}:{DRED}{e}{RESET}")
                resume_chunk_pos, attempt = self._handle_retry_attempt(
                    attempt, checkpoint_file
                )

        else:  # else block associated with while loop, executed if loop completes without break
            self._handle_synthesis_failure(max_retries)

    def _prepare_synthesis_environment(self, output_file, _tmp_folder_):
        """Prepares the environment for synthesis, handling temporary directory and checkpoint file."""
        out_dir = os.path.split(output_file)[0]
        _file_ = os.path.split(output_file)[1]

        if os.path.exists(os.path.join(out_dir, _tmp_folder_)):
            query = input(
                f"{DBLUE}Remove the {os.path.join(out_dir, _tmp_folder_)} directory (y/n)?{RESET} "
            ).lower() in ("y", "yes")
            if query:
                shutil.rmtree(os.path.join(out_dir, _tmp_folder_))

        _full_output_path_ = os.path.join(out_dir, _tmp_folder_, _file_)
        _check_file = os.path.splitext(_file_)[0] + ".ch"
        _check_dir_path_ = os.path.join(out_dir, _tmp_folder_, _check_file)
        checkpoint_file = _check_dir_path_
        start_chunk = 0

        if os.path.exists(checkpoint_file):
            logger.info(f"{DYELLOW}Found a Checkpoint file{RESET}")
            start_chunk = self._load_checkpoint(checkpoint_file)
            resume_chunk_pos = start_chunk * 1_000
            logger.info(f"{DYELLOW}Resuming from chunk{DBLUE} {start_chunk}{RESET}")
        else:
            resume_chunk_pos = start_chunk

        if not os.path.exists(_tmp_folder_):
            logger.info(
                f"{DYELLOW}Create temporary directory = {DBLUE}{_tmp_folder_}{RESET}"
            )
            os.makedirs(
                _tmp_folder_, exist_ok=True
            )  # Use makedirs to avoid FileExistsError

        logger.info(f"{DYELLOW}Start conversion{RESET}")
        return (
            out_dir,
            _file_,
            _full_output_path_,
            checkpoint_file,
            start_chunk,
            resume_chunk_pos,
        )

    def _get_output_filename_for_chunk(self, _full_output_path_, counter, start_chunk):
        """Generates the output filename for a chunk, resolving potential conflicts."""
        output_filename_base = f"{_full_output_path_}_{counter}.ogg"
        if os.path.exists(output_filename_base):
            if counter == start_chunk:
                print(f"{CYAN}Chunk vs file conflict: {BLUE}Resolving{RESET}")
                os.remove(output_filename_base)
            elif os.path.getsize(output_filename_base) != 0:  # Remove only if not empty
                os.remove(output_filename_base)
            return output_filename_base
        else:
            return output_filename_base

    def _synthesize_chunk_and_save(self, chunk, output_filename):
        """Synthesizes a text chunk to speech and saves it to a file."""
        tts = gTTS(text=chunk, lang="en", slow=False)
        tts.save(output_filename)

    def _save_checkpoint(self, checkpoint_file, counter):
        """Saves the current processing chunk counter to the checkpoint file."""
        with open(checkpoint_file, "w") as f:
            f.write(str(counter))

    def _load_checkpoint(self, checkpoint_file):
        """Loads the checkpoint from file and returns the chunk counter."""
        with open(checkpoint_file, "r") as f:
            return int(f.read())

    def _handle_retry_attempt(self, attempt, checkpoint_file):
        """Handles retry logic with exponential backoff and returns updated attempt and resume position."""
        attempt += 1
        for _sec_ in range(2**attempt, 0, -1):
            print(f"{BWHITE}Resume in {DBLUE}{_sec_}{RESET}", end="\r")
            time.sleep(1)  # Add a sleep to actually wait

        resume_chunk_pos = 0
        if os.path.exists(checkpoint_file):
            resume_chunk_pos = self._load_checkpoint(checkpoint_file) * 1_000
        return resume_chunk_pos, attempt

    def _handle_synthesis_success(self, _tmp_folder_, output_file, num_temp_files):
        """Handles successful synthesis, joins audio files and cleans up."""
        print(
            f"{FMAGENTA}Conversion success✅. \n  {FCYAN}INFO\t Create masterfile{RESET}"
        )
        if num_temp_files > 2:  # More robust check for multiple files
            from .JoinAudios import (
                JoinAudios,
            )  # Assuming JoinAudios is in the same directory or in PYTHONPATH

            joiner = JoinAudios(_tmp_folder_, masterfile=output_file)
            joiner.worker()

    def _handle_synthesis_failure(self, max_retries):
        """Handles synthesis failure after maximum retries."""
        print(
            f"{RED}Maximum retries reached. Unable to complete the operation after {DMAGENTA} {max_retries} attempts.{RESET}"
        )
        sys.exit(2)

    @staticmethod
    def pdf_to_text(pdf_path):
        logger.info(f"{GREEN} Initializing pdf to text conversion{RESET}")
        try:
            with open(pdf_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text = ""
                _pg_ = 0
                print(f"{YELLOW}Convert pages..{RESET}")
                for page_num in range(len(pdf_reader.pages)):
                    _pg_ += 1
                    logger.info(f"Page {DBLUE}{_pg_}{RESET}/{len(pdf_reader.pages)}")
                    page = pdf_reader.pages[page_num]
                    text += page.extract_text()
                print(f"{DGREEN}Ok{RESET}")
                return text
        except Exception as e:
            logger.error(
                f"{DRED}Failed to extract text from '{YELLOW}{pdf_path}'{RESET}:\n {e}"
            )

    @staticmethod
    def text_file(input_file):
        try:
            with open(input_file, "r", errors="ignore") as file:
                text = file.read().replace("\n", " ")
            return text
        except FileNotFoundError:
            logger.error("File '{}' was not found.📁".format(input_file))
        except Exception as e:
            logger.error(f"{DRED}{str(e)}{RESET}")

    @staticmethod
    def docx_to_text(docx_path):
        try:
            logger.info(f"{BLUE} Converting {docx_path} to text{RESET}")
            doc = Document(docx_path)
            paragraphs = [paragraph.text for paragraph in doc.paragraphs]
            return "\n".join(paragraphs)
        except FileNotFoundError:
            logger.error(f"File '{docx_path}' was not found.📁")
        except Exception as e:
            logger.error(f"{DRED}Error converting {docx_path} to text: {e}{RESET}")

    def audiofy(self):
        """
        Handle input files based on type to initialize conversion sequence
        """
        input_list = self.preprocess()
        ls = ("pdf", "docx", "doc", "txt", "ppt", "pptx")
        input_list = [item for item in input_list if item.lower().endswith(tuple(ls))]

        for input_file in input_list:
            output_file = os.path.split(input_file)[-1].split(".")[0] + ".ogg"

            if input_file.endswith(".pdf"):
                text = self.pdf_to_text(input_file)

            elif input_file.lower().endswith(tuple(_ext_word)):
                text = self.docx_to_text(input_file)

            elif input_file.endswith(".txt"):
                text = self.text_file(input_file)

            elif input_file.split(".")[-1] in ("ppt", "pptx"):
                # conv = FBot()
                # word = conv.ppt_to_word() # This seems to be converting ppt to word, not text, might be incorrect for audiofy
                text = (
                    self.docx_to_text()
                )  # This is also likely incorrect, docx_to_text expects docx_path

            else:
                logger.error(
                    "Unsupported file format. Please provide a PDF, txt, or Word document."
                )
                sys.exit(1)

            try:
                self.Synthesise(text, output_file)  # Pass the extracted text here
            except KeyboardInterrupt:
                print("\nQuit❕")
                sys.exit(1)

    def preprocess(self):  # Dummy preprocess, replace with actual logic if needed.
        return ["input.txt"]  # Example, adapt to your actual input sourc


class VideoConverter:
    def __init__(self, callback, **kwargs):
        self.callback = callback

    def preprocess(self):
        files_to_process = []

        if os.path.isfile(self.input_file):
            files_to_process.append(self.input_file)
        elif os.path.isdir(self.input_file):
            if os.listdir(self.input_file) is None:
                print(f"{RED}Cannot work with empty folder{RESET}")
                sys.exit(1)
            for file in os.listdir(self.input_file):
                file_path = os.path.join(self.input_file, file)
                if os.path.isfile(file_path):
                    files_to_process.append(file_path)

        return files_to_process

    def CONVERT_VIDEO(self):
        try:
            input_list = self.preprocess()
            out_f = self.out_format.upper()
            input_list = [
                item
                for item in input_list
                if any(item.upper().endswith(ext) for ext in SUPPORTED_VIDEO_FORMATS)
            ]
            print(f"{DYELLOW}Initializing conversion..{RESET}")

            for file in tqdm(input_list):
                if out_f.upper() in Video_codecs.keys():
                    _, ext = os.path.splitext(file)
                    output_filename = _ + "." + out_f.lower()
                    # print(output_filename)
                elif (
                    out_f.upper() in SUPPORTED_VIDEO_FORMATS
                    and not out_f.upper() in Video_codecs.keys()
                ):
                    print(
                        f"{RED}Unsupported output format --> Pending Implementation{RESET}"
                    )
                    sys.exit(1)
                else:
                    print(f"{RED}Unsupported output format{RESET}")
                    sys.exit(1)

                """Load the video file"""
                print(f"{DBLUE}Load file{RESET}")
                video = VideoFileClip(file)
                """Export the video to a different format"""
                print(f"{DMAGENTA}Converting file to {output_filename}{RESET}")
                video.write_videofile(output_filename, codec=Video_codecs[out_f])
                """Close the video file"""
                print(f"{DGREEN}Done{RESET}")
                video.close()
        except KeyboardInterrupt:
            print("\nQuit❕")
            sys.exit(1)
        except Exception as e:
            print(e)


class AudioConverter:
    def __init__(self, callback, **kwargs):
        self.callback = callback

    def preprocess(self):
        files_to_process = []

        if os.path.isfile(self.input_file):
            files_to_process.append(self.input_file)
        elif os.path.isdir(self.input_file):
            if os.listdir(self.input_file) is None:
                print(f"{RED}Cannot work with empty folder{RESET}")
                sys.exit(1)
            for file in os.listdir(self.input_file):
                file_path = os.path.join(self.input_file, file)
                if os.path.isfile(file_path):
                    files_to_process.append(file_path)

        return files_to_process

    def pydub_conv(self, obj, out_format):
        try:
            input_list = [obj]
            out_f = out_format
            input_list = [
                item
                for item in input_list
                if any(
                    item.lower().endswith(ext) for ext in SUPPORTED_AUDIO_FORMATS_INPUT
                )
            ]
            print(f"{DYELLOW}Initializing conversion..{RESET}")

            def wav_redudancy(file):
                # Load the mp3 file using Pydub
                audio = pydub.AudioSegment.from_file(file, fmt)
                # Export the audio to a temporary file in wav format (ffmpeg can convert from wav to m4a)
                audio.export("temp.wav", format="wav")

            for file in tqdm(input_list):
                if out_f.lower() in SUPPORTED_AUDIO_FORMATS_OUTPUT:
                    _, ext = os.path.splitext(file)
                    output_filename = _ + "." + out_f
                    fmt = ext[1:]
                    # print(fmt, out_f)
                    audio = pydub.AudioSegment.from_file(file, fmt)
                    print(f"{DMAGENTA}Converting to {output_filename}{RESET}")
                    audio.export(output_filename, format=out_f)
                    # new_audio = pydub.AudioSegment.from_file('output_audio.')
                    print(f"{DGREEN}Done{RESET}")

                elif file[-3:].lower() == "m4a" or out_f.lower() == "m4a":
                    _m4a_main_(file, out_f)

                elif (
                    out_f.lower() in SUPPORTED_AUDIO_FORMATS_INPUT
                    and not SUPPORTED_AUDIO_FORMATS_OUTPUT
                ):
                    print("Pending Implemantation For the format")

                else:
                    print(f"{RED}Unsupported output format{RESET}")
                    sys.exit(1)

        except KeyboardInterrupt:
            print("\nQuit❕")
            sys.exit(1)
        except Exception as e:
            print(f"{RED}{e}{RED}")


class ImageConverter:
    def __init__(self, callback, **kwargs):
        self.callback = callback

    def preprocess(self):
        try:
            files_to_process = []

            if os.path.isfile(self.input_file):
                files_to_process.append(self.input_file)
            elif os.path.isdir(self.input_file):
                if os.listdir(self.input_file) is None:
                    print("Cannot work with empty folder")
                    sys.exit(1)
                for file in os.listdir(self.input_file):
                    file_path = os.path.join(self.input_file, file)
                    if os.path.isfile(file_path):
                        files_to_process.append(file_path)

            return files_to_process
        except FileNotFoundError:
            print("File not found❕")
            sys.exit(1)

    def convert_image(self, file):
        try:
            input_list = self.preprocess()
            out_f = self.out_format.upper()
            input_list = [
                item
                for item in input_list
                if any(
                    item.lower().endswith(ext)
                    for ext in SUPPORTED_IMAGE_FORMATS.values()
                )
            ]

            for file in tqdm(input_list):
                if out_f.upper() in SUPPORTED_IMAGE_FORMATS:
                    _ = os.path.splitext(file)[0]
                    output_filename = _ + SUPPORTED_IMAGE_FORMATS[out_f].lower()
                else:
                    print("Unsupported output format")
                    sys.exit(1)
                """Load the image using OpenCV: """
                print(f"{DYELLOW}Reading input image..{RESET}")
                img = cv2.imread(file)
                """Convert the OpenCV image to a PIL image: """
                print(f"{DMAGENTA}Converting to PIL image{RESET}")
                pil_img = Image.fromarray(cv2.cvtColor(img, cv2.COLOR_BGR2RGB))
                """Save the PIL image to a different format: """
                print(f"\033[1;36mSaving image as {output_filename}{RESET}")
                pil_img.save(output_filename, out_f)
                print(f"{DGREEN}Done ✅{RESET}")
                """
                Load the image back into OpenCV:
                # print(f"{DMAGENTA}Load and display image{RESET}")
                # opencv_img = cv2.imread(output_filename)
                Display the images:
                # cv2.imshow('OpenCV Image', opencv_img)
                # opencv_img.show()
                Wait for the user to press a key and close the windows:
                # cv2.waitKey(0)
                # cv2.destroyAllWindows()
                """
        except KeyboardInterrupt:
            print("\nQuit❕")
            sys.exit(1)
        except AssertionError:
            print("Assertion failed.")
        except KeyError:
            print(
                f"{RED}ERROR:\tPending Implementation for{ICYAN} {out_f} {BWHITE}format{RESET}"
            )
        except Exception as e:
            print(f"{RED}{e}{RESET}")
